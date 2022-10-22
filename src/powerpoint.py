import collections.abc
from operator import imod # インポートしないとエラーが発生する
import pptx
from pptx.enum.shapes import MSO_SHAPE    # 図形の定義がされているクラス
from pptx.dml.color import RGBColor
from pptx.util import Cm,Pt,Inches              # 単位指定をするクラス(センチメートル, ポイント単位)
import os
from PIL import Image 
from pptx.enum.text import PP_ALIGN       # 段落の水平位置のEnume

def put_pic(path, pic_left, pic_top, pic_width, pic_height):
    #画像追加
    image = slide.shapes.add_picture(path, pic_left, pic_top, pic_width, pic_height) 
    image.line.color.rgb = RGBColor(0, 0, 0)
    image.line.width = Pt(1.5)

def put_text(pic_left, pic_top, str, size):
    #テキストを追加
    textbox = slide.shapes.add_textbox(pic_left, pic_top, Pt(size), Pt(size))
    tf = textbox.text_frame
    tf.text = str
    tf.paragraphs[0].font.size = Pt(size)  # font size
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT

def put_arrow(pic_left2, pic_top2, size):
    #矢印出力
    rect0 = slide.shapes.add_shape(		# shapeオブジェクト➀を追加
            MSO_SHAPE.RIGHT_ARROW,   	                    # 図形の種類を[丸角四角形]に指定
            pic_left2, pic_top2,               # 挿入位置の指定：左からの座標と上からの座標の指定
            size, size)               # 挿入図形の幅と高さの指定
    rect0.fill.solid()                                   # shapeオブジェクト➀を単色で塗り潰す
    rect0.fill.fore_color.rgb = RGBColor(74, 126, 187)  # RGB指定で色を指定

def put_sign(SIGN_DIR, sign_names):
     pic_top3 = Cm(0)
     width = Cm(5.49)
     height = Cm(3.94)
     for i,name in enumerate(sign_names):
         if(i < 2):
            slide.shapes.add_picture(SIGN_DIR+name, Cm(-6.5), pic_top3, height=height) 
            pic_top3 += height
         elif i == 2:
            pic_top3 = Cm(-1.05)
            width = Cm(5.92)
            height = Cm(4.92)
            slide.shapes.add_picture(SIGN_DIR+name, Cm(29.69), pic_top3, height=height) 
            #pic_top3 += height
         elif i == 3:
            slide.shapes.add_picture(SIGN_DIR+name, Cm(31.69), Cm(2.8), height=height) 
            #pic_top3 += height
         elif i == 4:
            height = Cm(4.92)
            slide.shapes.add_picture(SIGN_DIR+name, Cm(30.96),Cm(5.9), height=height) 
         else:
            #width = Cm(4.92)
            height = Cm(5.61)
            slide.shapes.add_picture(SIGN_DIR+name, Cm(30.96), Cm(11.12), height=height) 
    


        
#画像の格納ディレクトリ
IMG_DIR = "./output_func5/"
#img画像のファイル名を取得
img_names = os.listdir(IMG_DIR)
img_names = [name for name in img_names if name.endswith(".jpeg")]
img_names.sort()#昇順にsort

#画像の格納ディレクトリ
SIGN_DIR = "./sign/"
#img画像のファイル名を取得
sign_names = os.listdir(SIGN_DIR)
sign_names = [name for name in sign_names if name.endswith(".png")]
sign_names.sort()#昇順にsort

prs = pptx.Presentation()
prs.slide_width = Inches(11.69) #A4サイズ)
prs.slide_height = Inches(8.27)
slide_width = prs.slide_width
slide_height = prs.slide_height

#画像のアスペクト比を取得
im = Image.open(IMG_DIR + img_names[0])
aspect_ratio = im.width /im.height

#画像の高さの設定と幅の取得
pic_height = Cm(9.5)
pic_width = aspect_ratio * pic_height

#タイトルスライド
# slide = prs.slides.add_slide(prs.slide_layouts[0])
# slide.shapes.title.text = "タイトルを入力"

slide = prs.slides.add_slide(prs.slide_layouts[6])
textbox = slide.shapes.add_textbox(slide_width/2-Inches(5)/2, slide_height/2-Inches(1)/2, Inches(5), Inches(1))
tf = textbox.text_frame
tf.text = "タイトルを入力"
tf.paragraphs[0].font.size = Pt(50)  # font size
tf.paragraphs[0].alignment = PP_ALIGN.CENTER


#画像を１枚のパワポに出力 1段4枚ずつ
pic_left = Cm(2.5)
pic_left2 = Cm(7.2)
for i, name in enumerate(img_names):
    path = IMG_DIR + name

    if i % 8 == 0:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        pic_top = Cm(0.5)

    #画像追加
    put_pic(path, pic_left, pic_top, pic_width, pic_height)
    #テキストを追加
    put_text(pic_left-Cm(1), pic_top, str(i+1), 28)

    pic_left += Cm(7)

    if i % 4 == 3:
        pic_top += Cm(10)
        pic_left = Cm(2.5)
        pic_left2 = Cm(7.2)
        pic_top2 += Cm(11)
    elif i != len(img_names)-1:
    #矢印出力
        pic_top2 = pic_top + pic_height/2 - Cm(1)
        put_arrow(pic_left2, pic_top2, Cm(2))
        pic_left2 += Cm(7) 

#画像を２枚ずつパワポに出力
pic_height = Cm(16)
pic_width = aspect_ratio * pic_height
pic_top = ( slide_height - pic_height ) / 2

# 連番で2枚ずつのスライドを作る疑似コード
pre_path = None
for i,name in enumerate(img_names):
    path = IMG_DIR + name
    if i==0:
        pre_path = path
        continue
    #スライドを増やす
    slide = prs.slides.add_slide(prs.slide_layouts[6]) 
    #pre_pathの画像を←に配置
    pic_left = ( slide_width/2 - pic_width ) / 2
    put_pic(pre_path, pic_left, pic_top, pic_width, pic_height)
    put_text(pic_left-Cm(1.5), pic_top, str(i), 36)
    pre_path = path
    #pathの画像を→に配置
    pic_left += slide_width/2
    put_pic(path, pic_left, pic_top, pic_width, pic_height)
    #矢印を追加
    ratio = 0.45
    pic_left2 = slide_width/2 - pic_width*ratio/2
    pic_top2 = slide_height/2 - pic_width*ratio/2
    put_arrow(pic_left2, pic_top2, pic_width*ratio)
    #手の写真を追加
    put_sign(SIGN_DIR, sign_names)
    #テキストを追加
    put_text(pic_left-Cm(1.5), pic_top, str(i+1), 36)
    pre_path = path
    
prs.save("./test.pptx")


    